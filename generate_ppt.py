import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_lms_presentation(output_path):
    prs = Presentation()
    # Widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Palette
    NAVY = RGBColor(15, 23, 42)        # #0f172a
    BLUE_ACCENT = RGBColor(0, 102, 255) # #0066ff
    LIGHT_BG = RGBColor(248, 250, 252)  # #f8fafc
    CARD_BG = RGBColor(255, 255, 255)   # #ffffff
    TEXT_DARK = RGBColor(30, 41, 59)     # #1e293b
    TEXT_MUTED = RGBColor(100, 116, 139) # #64748b
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(245, 158, 11)      # #f59e0b
    GREEN = RGBColor(16, 185, 129)     # #10b981
    BORDER_COLOR = RGBColor(226, 232, 240)
    PURPLE = RGBColor(147, 51, 234)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, subtitle_text=""):
        header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = NAVY
        header_box.line.color.rgb = NAVY
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.55))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = GOLD
            p2.font.name = 'Calibri'

    def add_card(slide, left, top, width, height, title, items, title_color=BLUE_ACCENT):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        
        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"• {item}"
            p_item.font.size = Pt(11.5)
            p_item.font.color.rgb = TEXT_DARK
            p_item.font.name = 'Calibri'
            p_item.space_after = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.color.rgb = NAVY

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(0.15), Inches(3.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_ACCENT
    bar.line.color.rgb = BLUE_ACCENT

    txBox = slide1.shapes.add_textbox(Inches(1.6), Inches(2.1), Inches(10.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "LMS SYSTEM DATA FLOW & ARCHITECTURE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'

    p2 = tf.add_paragraph()
    p2.text = "End-to-End Table Data Flow, ER Connections & 22-Table Relational Model for Students"
    p2.font.size = Pt(19)
    p2.font.color.rgb = GOLD
    p2.font.name = 'Calibri'
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "Target Audience: Data Analytics & Engineering Students | System: AppTechno LMS"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.font.name = 'Calibri'
    p3.space_before = Pt(24)

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Executive Summary & EdTech Data Analytics", "Why Data Analytics Students Need to Understand LMS Architecture")

    add_card(slide2, 0.8, 1.6, 5.6, 5.2, "🎯 Purpose of an Enterprise LMS", [
        "Automates the complete student lifecycle from initial inquiry to job placement.",
        "Generates multi-dimensional business data across marketing, finance, academics, and recruitment.",
        "Provides real-world datasets ideal for building Business Intelligence (BI) dashboards.",
        "Requires scalable relational database schemas to ensure transactional consistency."
    ], BLUE_ACCENT)

    add_card(slide2, 6.8, 1.6, 5.6, 5.2, "📈 Key Data Analytics Domain Scope", [
        "Marketing & Sales Funnel: Prospect conversion & marketing ROI.",
        "Financial Analytics: Revenue realization & fee collection efficiency.",
        "Academic & Engagement: Attendance tracking & churn prediction.",
        "Evaluation Analytics: Exam performance & weekly trainer feedback.",
        "Career Analytics: Mock interview readiness & placement conversion."
    ], GOLD)

    # -------------------------------------------------------------
    # SLIDE 3: NEW! COMPLETE STEP-BY-STEP TABLE DATA FLOW DIAGRAM
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Complete Step-by-Step Table Data Flow Diagram", "How Data Flows sequentially through Database Tables in 6 Business Steps")

    flow_steps = [
        ("1. INQUIRY", "leads\nlead_activities", "Prospect registers interest. Calls & emails logged.", BLUE_ACCENT),
        ("2. ADMISSION", "users\nregistrations", "Lead converts to User. Fee payment recorded.", GOLD),
        ("3. COHORT", "courses, batches\nbatch_students", "Student assigned to course batch cohort.", GREEN),
        ("4. DAILY OPS", "attendance\ntime_tracking", "Daily QR login hours & portal time logged.", BLUE_ACCENT),
        ("5. LEARNING", "projects, tasks\nassessments", "Tasks submitted, exams scored, feedback recorded.", PURPLE),
        ("6. CAREER", "mock_interviews\njob_applications", "Mock scores logged & job selection tracked.", GOLD)
    ]

    card_w = 1.75
    gap_w = 0.2
    start_x = 0.8
    for idx, (step_num, tables, desc, col) in enumerate(flow_steps):
        x_pos = start_x + idx * (card_w + gap_w)
        
        # Main step box
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.6), Inches(card_w), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = col
        card.line.width = Pt(2.0)
        
        txBox = slide3.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.75), Inches(card_w - 0.2), Inches(4.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Step Title
        p = tf.paragraphs[0]
        p.text = step_num
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        
        # Table Names Badge
        p_tbl = tf.add_paragraph()
        p_tbl.text = f"\n[ Tables ]\n{tables}"
        p_tbl.font.size = Pt(11)
        p_tbl.font.bold = True
        p_tbl.font.color.rgb = NAVY
        p_tbl.font.name = 'Calibri'
        p_tbl.alignment = PP_ALIGN.CENTER
        
        # Description
        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_DARK
        p_desc.font.name = 'Calibri'
        p_desc.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 4: NEW! TABLE-TO-TABLE TRANSITION DATA PIPELINE
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Table-to-Table Data Transition Pipeline", "Data Inputs & Outputs Passing Between Interconnected Database Tables")

    pipelines = [
        ("Input Stage 1: Lead Capture", "leads (leadId) ➔ lead_activities (activityId)", "Marketers capture lead details. Every call/email adds a row in lead_activities referencing leadId.", BLUE_ACCENT),
        ("Transition Stage 2: Enrollment", "leads ➔ users (userId) ➔ registrations", "When lead status changes to CONVERTED, a new record is created in users. Fees logged in registrations.", GOLD),
        ("Execution Stage 3: Training", "users ➔ batch_students ➔ attendance & tasks", "Student userId is mapped to batchId. Daily attendance logs & task progress link to userId & batchId.", GREEN),
        ("Evaluation Stage 4: Assessments", "users ➔ assessment_submissions & feedback", "Exams taken create records in assessment_submissions. Weekly ratings recorded in feedback table.", PURPLE),
        ("Outcome Stage 5: Job Hiring", "users ➔ mock_interviews ➔ job_applications", "Mock interviews prepare students. Job applications track status from APPLIED to SELECTED.", GOLD)
    ]

    top_pos = 1.5
    for stage, tables_flow, explanation, col in pipelines:
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.733), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.08), Inches(11.3), Inches(0.8))
        tf = txBox.text_frame
        
        p = tf.paragraphs[0]
        p.text = f"{stage}   |   {tables_flow}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = 'Calibri'
        
        p2 = tf.add_paragraph()
        p2.text = explanation
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_DARK
        p2.font.name = 'Calibri'

        top_pos += 1.1

    # -------------------------------------------------------------
    # SLIDE 5: Central Table Connection Architecture (ER Hubs)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. Central Table Connection Architecture", "The 3 Primary Hub Entities Connecting All 22 System Tables")

    add_card(slide5, 0.8, 1.6, 3.6, 5.2, "👤 HUB 1: 'users' Table\n(Central Entity)", [
        "admin_permissions (1:1)",
        "batch_students (1:N)",
        "registrations (1:N)",
        "documents (1:N)",
        "attendance (1:N)",
        "leave_requests (1:N)",
        "time_tracking (1:N)",
        "tasks (1:N)",
        "assessment_submissions (1:N)",
        "feedback (1:N)",
        "leads & lead_activities (1:N)",
        "job_applications (1:N)",
        "mock_interviews (1:N)",
        "communication_practice (1:N)",
        "notifications & messages (1:N)"
    ], BLUE_ACCENT)

    add_card(slide5, 4.86, 1.6, 3.6, 5.2, "👥 HUB 2: 'batches' Table\n(Academic Cohort Hub)", [
        "courses (N:1 -> courseId)",
        "users (N:1 -> trainerId)",
        "batch_students (1:N)",
        "registrations (1:N)",
        "attendance (1:N)",
        "leave_requests (1:N)",
        "projects (1:N)",
        "videos (1:N)",
        "feedback (1:N)"
    ], GREEN)

    add_card(slide5, 8.93, 1.6, 3.6, 5.2, "📚 HUB 3: 'courses' & Sub-Hubs\n(Catalog & Core Sub-Entities)", [
        "courses -> batches (1:N)",
        "courses -> registrations (1:N)",
        "courses -> assessments (1:N)",
        "projects -> tasks (1:N)",
        "leads -> lead_activities (1:N)",
        "assessments -> submissions (1:N)",
        "jobs -> job_applications (1:N)"
    ], GOLD)

    # -------------------------------------------------------------
    # SLIDE 6: Complete Foreign Key Connection Map
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Complete Foreign Key (FK) Connection Map", "Explicit Mapping of All Foreign Keys Across All 22 Database Tables")

    rows, cols = 11, 3
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4)
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(5.433)

    headers = ["Child Table (FK Owner)", "Foreign Key (FK Column)", "Referenced Parent Table & Key"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

    fk_mappings = [
        ("admin_permissions", "userId", "users.id (1:1 cascade delete)"),
        ("batches", "courseId, trainerId", "courses.id, users.id (Trainer)"),
        ("batch_students", "batchId, studentId", "batches.id, users.id (Enrollment junction)"),
        ("leads", "assignedToId, createdById", "users.id (Marketer / Creator)"),
        ("lead_activities", "leadId, userId", "leads.id, users.id (Touchpoint logs)"),
        ("registrations", "studentId, courseId, batchId", "users.id, courses.id, batches.id"),
        ("attendance & leaves", "studentId/userId, batchId", "users.id, batches.id (Daily ops & approvals)"),
        ("tasks", "projectId, studentId", "projects.id, users.id (Task ticketing)"),
        ("assessment_submissions", "assessmentId, studentId", "assessments.id, users.id (Exam scores)"),
        ("job_applications", "jobId, studentId", "jobs.id, users.id (Recruitment funnel)")
    ]

    for row_idx, data in enumerate(fk_mappings, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 1 else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # SLIDE 7: Database Schema (22 Tables Breakdown - Part 1)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Relational Database Schema (Part 1: Core, Marketing, Finance)", "22 Tables Grouped into Functional Modules")

    add_card(slide7, 0.8, 1.6, 5.6, 2.5, "Module A: Users & RBAC", [
        "users: Account master (email, role, phone, studentId)",
        "admin_permissions: Granular administrative flags"
    ], BLUE_ACCENT)

    add_card(slide7, 6.8, 1.6, 5.6, 2.5, "Module B: Courses & Batches", [
        "courses: Master course catalog & base tuition fees",
        "batches: Academic cohorts, schedules & trainer assignment",
        "batch_students: Junction table mapping student enrollments"
    ], GOLD)

    add_card(slide7, 0.8, 4.3, 5.6, 2.5, "Module C: Marketing Funnel", [
        "leads: Prospect pipeline (source, status, assignee)",
        "lead_activities: Sales touchpoint log (calls, emails)"
    ], GREEN)

    add_card(slide7, 6.8, 4.3, 5.6, 2.5, "Module D: Finance & Admissions", [
        "registrations: Fee tracking (feeAmount, feePaid, receipt)",
        "documents: Student KYC uploads & admin verification status"
    ], BLUE_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 8: Database Schema (22 Tables Breakdown - Part 2)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Relational Database Schema (Part 2: Ops, Academics, Career)", "Attendance, Evaluation & Placement Modules")

    add_card(slide8, 0.8, 1.6, 5.6, 2.5, "Module E: Attendance & Time Logs", [
        "attendance: Daily class QR/Geo attendance records",
        "leave_requests: Leave applications, proof URLs & status",
        "time_tracking: Lab portal active duration tracking"
    ], BLUE_ACCENT)

    add_card(slide8, 6.8, 1.6, 5.6, 2.5, "Module F: Projects & Content", [
        "projects: Capstone & module projects per batch",
        "tasks: Task tickets per student & plagiarism flags",
        "videos: Lecture recordings & timeline chapters"
    ], GOLD)

    add_card(slide8, 0.8, 4.3, 5.6, 2.5, "Module G: Evaluation Engine", [
        "assessments: Exam master (Coding & Aptitude tests)",
        "assessment_submissions: Student test scores & answers",
        "feedback: Weekly ratings for batch & trainer performance"
    ], GREEN)

    add_card(slide8, 6.8, 4.3, 5.6, 2.5, "Module H: Placement & Career", [
        "jobs: Corporate hiring drives & salary packages",
        "job_applications: Recruitment funnel stage tracking",
        "mock_interviews & communication_practice: Drill scores"
    ], BLUE_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 9: Analytics Project 1 - Marketing Funnel
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Data Analytics Project 1: Lead Conversion Funnel", "Evaluating Marketing Channels & Sales Productivity")

    add_card(slide9, 0.8, 1.6, 5.6, 5.2, "📋 Project Specifications", [
        "Primary Tables: leads, lead_activities, registrations",
        "Goal: Measure marketing ROI and sales conversion efficiency.",
        "Analytical Focus: Tracking lead journeys from NEW to CONVERTED.",
        "Business Impact: Optimizing marketing spend on high-performing channels."
    ], BLUE_ACCENT)

    add_card(slide9, 6.8, 1.6, 5.6, 5.2, "💡 Key KPIs & Analytics Questions", [
        "1. Lead Conversion Rate (LCR) by Source (Google Ads vs LinkedIn vs Referral).",
        "2. Marketer Sales Productivity: Conversion % per assigned sales rep.",
        "3. Average Touchpoints: Number of calls/emails needed before student registration.",
        "4. Drop-off Analysis: High-loss stages in the marketing funnel."
    ], GOLD)

    # -------------------------------------------------------------
    # SLIDE 10: Analytics Project 2 - Student Churn Risk
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Data Analytics Project 2: Student Churn Prediction", "Predictive Analytics for Academic Retention & Early Intervention")

    add_card(slide10, 0.8, 1.6, 5.6, 5.2, "📊 Metric Formulas", [
        "Attendance Percentage = (Count of PRESENT / Total Batch Days) * 100",
        "Task Overdue Rate = (Count of OVERDUE Tasks / Total Assigned Tasks) * 100",
        "Primary Tables: attendance, tasks, leave_requests, feedback",
        "Goal: Identify at-risk students before course abandonment."
    ], BLUE_ACCENT)

    add_card(slide10, 6.8, 1.6, 5.6, 5.2, "🎯 Churn Risk Rules & Actions", [
        "High Risk Threshold: Attendance < 75% AND Overdue Tasks > 30%.",
        "Early Warning Indicator: Consistently low weekly feedback rating (< 3 stars).",
        "Intervention Trigger: Automated alert to batch trainers and academic advisors.",
        "Outcome: Increased course completion rate and student success."
    ], GREEN)

    # -------------------------------------------------------------
    # SLIDE 11: Analytics Project 3 - Financial Realization
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "10. Data Analytics Project 3: Revenue & Fee Realization", "Financial Analytics & Tuition Collection Efficiency")

    add_card(slide11, 0.8, 1.6, 5.6, 5.2, "💰 Revenue Metrics", [
        "Total Agreed Tuition Revenue = SUM(feeAmount)",
        "Collected Cash Realization = SUM(feePaid)",
        "Outstanding Fee Balance = SUM(feeAmount - feePaid)",
        "Collection Efficiency % = (Total Fee Paid / Total Fee Amount) * 100",
        "Primary Tables: registrations, courses, batches"
    ], GOLD)

    add_card(slide11, 6.8, 1.6, 5.6, 5.2, "📈 BI Dashboard Questions", [
        "1. Which course generates the highest gross revenue per batch?",
        "2. What is the total uncollected fee balance across active batches?",
        "3. Fee payment delinquency trends by student enrollment month.",
        "4. Revenue realization forecasting for upcoming batch end-dates."
    ], BLUE_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 12: Analytics Projects 4 & 5 - Quality & Placement
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "11. Analytics Projects 4 & 5: Quality & Career Funnel", "Faculty Evaluation & Job Placement Readiness Analytics")

    add_card(slide12, 0.8, 1.6, 5.6, 5.2, "⭐ Project 4: Trainer Satisfaction Index", [
        "Primary Tables: feedback, batches, users",
        "Metric: Weekly 1 to 5 star rating average per trainer.",
        "NLP Sentiment Analysis: Text mining qualitative comments in feedback.comments.",
        "Goal: Ensuring consistent teaching quality across cohorts."
    ], GREEN)

    add_card(slide12, 6.8, 1.6, 5.6, 5.2, "💼 Project 5: Placement Funnel Performance", [
        "Primary Tables: job_applications, jobs, mock_interviews",
        "Metric: Correlation between Mock Score (>80) and Job Selection.",
        "Recruitment Funnel Conversion %: APPLIED -> SHORTLISTED -> INTERVIEW -> SELECTED.",
        "Goal: Maximizing corporate hiring rates for students."
    ], GOLD)

    # -------------------------------------------------------------
    # SLIDE 13: Summary & Takeaways
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "12. Summary & Key Takeaways for Students", "Connecting Database Architecture with Real-World BI & Analytics")

    add_card(slide13, 0.8, 1.6, 11.733, 5.2, "🎓 Essential Lessons for Data Analytics Students", [
        "1. Sequential Data Flow: Data moves in 6 distinct steps from Inquiry to Job Placement.",
        "2. Table Transitions: Foreign keys (leadId, userId, batchId) connect tables sequentially.",
        "3. Central Hub Pattern: 'users', 'batches', and 'courses' form the core backbone of all sub-modules.",
        "4. Actionable Insights: Raw transactional tables (logs, marks, fees) drive strategic business decisions.",
        "5. Complete Presentation & Reference guide available in workspace file: LMS_DATA_ARCHITECTURE_AND_FLOW.md"
    ], BLUE_ACCENT)

    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    out_file = sys.argv[1] if len(sys.argv) > 1 else "LMS_Data_Architecture_and_Flow.pptx"
    create_lms_presentation(out_file)
